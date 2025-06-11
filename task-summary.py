# Required libraries
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
import os
from datetime import datetime
import streamlit as st

# Step 1: Load Excel File
def load_status_report(file):
    return pd.read_excel(file, sheet_name=None)

# Step 2: Create Task Summary Sheet
def create_task_summary(df):
    summary = df.groupby(['Team Member', 'Status']).agg({'Task': 'count', 'Hours': 'sum'}).reset_index()
    summary_pivot = summary.pivot(index='Team Member', columns='Status', values='Task').fillna(0)
    summary_pivot['Total Hours'] = df.groupby('Team Member')['Hours'].sum()
    return summary_pivot.reset_index()

# Step 3: Load Incident Report (Optional)
def load_incident_report(file):
    return pd.read_excel(file)

# Step 4: Generate Dashboard Charts
def generate_dashboard(summary_df):
    charts_dir = "charts"
    os.makedirs(charts_dir, exist_ok=True)

    # Bar chart - tasks per status
    status_counts = summary_df.set_index('Team Member').drop(columns='Total Hours').sum()
    status_counts.plot(kind='bar', title='Task Status Summary')
    plt.ylabel('Task Count')
    plt.tight_layout()
    status_chart_path = os.path.join(charts_dir, 'status_chart.png')
    plt.savefig(status_chart_path)
    plt.close()

    # Pie chart - total hours per person
    summary_df.set_index('Team Member')['Total Hours'].plot(kind='pie', autopct='%1.1f%%', title='Total Hours by Member')
    plt.ylabel('')
    plt.tight_layout()
    hours_chart_path = os.path.join(charts_dir, 'hours_chart.png')
    plt.savefig(hours_chart_path)
    plt.close()

    return status_chart_path, hours_chart_path

# Step 5: Generate PPT Report
def generate_ppt(summary_df, status_chart, hours_chart, output_file):
    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Daily Status Summary"
    slide.placeholders[1].text = f"Generated on {datetime.now().strftime('%Y-%m-%d')}"

    # Summary Table slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Task Summary Table"
    rows, cols = summary_df.shape
    table = slide.shapes.add_table(rows+1, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table

    for col in range(cols):
        table.cell(0, col).text = summary_df.columns[col]
    for row in range(rows):
        for col in range(cols):
            table.cell(row+1, col).text = str(summary_df.iloc[row, col])

    # Chart slides
    for chart_path, title in zip([status_chart, hours_chart], ["Task Status", "Hours Distribution"]):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title
        slide.shapes.add_picture(chart_path, Inches(1), Inches(1.5), width=Inches(8))

    prs.save(output_file)

# === Streamlit Web App ===
st.title("📊 Project Report Generator Chatbot")

status_file = st.file_uploader("Upload Daily Status Excel File", type=[".xlsx"])
incident_file = st.file_uploader("(Optional) Upload Incident Report", type=[".xlsx"])

if status_file is not None:
    with st.spinner('Processing your report...'):
        data = load_status_report(status_file)
        task_df = data[list(data.keys())[0]]
        summary = create_task_summary(task_df)
        status_chart, hours_chart = generate_dashboard(summary)

        ppt_output = "Daily_Status_Report.pptx"
        generate_ppt(summary, status_chart, hours_chart, ppt_output)

        st.success("Report Generated Successfully!")

        st.download_button("📥 Download PowerPoint Report", open(ppt_output, "rb"), file_name=ppt_output)
        st.dataframe(summary)

        st.image(status_chart, caption="Task Status Chart")
        st.image(hours_chart, caption="Hours Distribution Chart")
